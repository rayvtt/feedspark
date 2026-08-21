#!/usr/bin/env python3
"""Convert a FeedSpark HTML deck into a themed, editable 16:9 PowerPoint deck.

This does NOT draw slides. It *populates a PowerPoint template*
(tools/templates/feedspark_deck.pptx), and that distinction is the whole point.

Every visual element -- card panels, accent bars, the decorative circles, the
gradient footer bar, the wordmark -- lives on that template's 18 named layouts.
Slides carry text in placeholders and nothing else. The result is a real
PowerPoint deck the client can restyle, re-theme, re-order or drop onto a
different layout and have it re-flow, rather than a picture of a deck welded
out of absolutely-positioned rectangles.

The HTML decks are continuous-scroll documents with variable-height sections;
PPTX is a fixed 13.333in x 7.5in canvas with no reflow. So the tool re-flows
semantically: it parses the deck's own component vocabulary (.stats/.card/
.tbl-wrap/.bars/.sc-grid/.tiers/.road/.callout/.note/.agenda/.contacts), maps
each onto the layout that fits it, and splits across as many slides as the copy
needs. Because every FeedSpark deck is built from the same component library,
this works on any of them (Superdry, Reiss, YuMOVE, Monsoon, ...).

Two rules keep the output client-ready:
  * Never add a shape. If something has no layout, add a layout to the template.
  * Never shrink text below MIN_OK to make it fit -- re-lay it out instead
    (fewer cards per slide, more slides). 7pt type is not a fit.

    pip install python-pptx pillow lxml
    python tools/deck_to_pptx.py docs/Superdry_Strategy_Review_AllTime.html out.pptx --audit

--audit is the QA loop: it reports the layouts used, anything shrunk, and
anything still over capacity. Ship when "still over capacity: 0".
"""
import re, sys, os, argparse
from lxml import html as LH
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import ImageFont

# ------------------------------------------------------------------ status ink
# Colour is inherited from the template theme for everything except short
# status-pill cells in tables, where the pill's meaning IS the colour.
MUTED      = RGBColor(0x76, 0x76, 0x76)
ORANGE_DP  = RGBColor(0xED, 0x6F, 0x0B)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Pillow metrics mirror preview_tmpl.py (Liberation Sans ~ Arial metrics), so
# heights measured here match what the QA previewer will measure.
LIB = "/usr/share/fonts/truetype/liberation/"
# Must match preview_tmpl.py's SC exactly: it rasterises at 120 px/inch, and the
# integer rounding of the pixel font size differs enough between 96 and 120 that
# a borderline label measured as one line here renders as two there.
SC = 120.0
# preview_tmpl.draw_textframe() wraps at `max(4, shape_width_px - 4)`; mirror that
# exactly rather than approximating, or borderline labels measure 1 line and render 2.
PAD_PX = 4.0
_fc = {}
def _font(pt_size, bold=False):
    key = (round(pt_size, 1), bold)
    if key in _fc: return _fc[key]
    path = LIB + ("LiberationSans-Bold.ttf" if bold else "LiberationSans-Regular.ttf")
    _fc[key] = ImageFont.truetype(path, max(6, int(round(pt_size * SC / 72.0))))
    return _fc[key]

def wrap_lines(text, pt_size, width_in, bold=False):
    """Wrap `text` to `width_in` inches at `pt_size`; return the line list."""
    f = _font(pt_size, bold)
    limit = max(4.0, width_in * SC - PAD_PX)
    out = []
    for para in (text or "").split("\n"):
        words, cur = para.split(), ""
        if not words:
            out.append("")
            continue
        for w in words:
            trial = (cur + " " + w).strip()
            if f.getlength(trial) <= limit or not cur:
                cur = trial
            else:
                out.append(cur); cur = w
        if cur: out.append(cur)
    return out

# ---------------------------------------------------------------- html parsing
# Liberation Sans (and Lato) carry none of these; PowerPoint would have to
# font-fall-back per glyph, which is inconsistent across Win/Mac/Slides. Status
# meaning is carried by colour in the renderer instead -- see status_col().
GLYPH_SUB = {"\u2713": "", "\u2691": "", "\u2715": "", "\u25c6": "", "\u2726": ""}

def safe_glyphs(s):
    for k, v in GLYPH_SUB.items():
        s = s.replace(k, v)
    return re.sub(r"\s{2,}", " ", s).strip()

def norm(s):
    """Collapse whitespace AND drop glyphs no common PowerPoint font carries."""
    return safe_glyphs(re.sub(r"\s+", " ", (s or "")).replace("\xa0", " ").strip())

def cls(el):
    return (el.get("class") or "")

def has(el, name):
    return name in cls(el).split()

def first(el, xp):
    """First xpath match or None -- .find() only speaks limited ElementPath."""
    if el is None: return None
    r = el.xpath(xp)
    return r[0] if r else None

def txt(el, strip_chk=True):
    """Visible text of an element, optionally dropping the `?` data-check badges."""
    if el is None: return ""
    c = LH.fromstring(LH.tostring(el))
    if strip_chk:
        for b in c.xpath("//span[contains(@class,'chk')]"):
            p = b.getparent()
            if p is not None: p.remove(b)
    for br in c.xpath("//br"):            # <br> is a real line break, not a space
        br.tail = "\ue000" + (br.tail or "")
    parts = [safe_glyphs(norm(x)) for x in c.text_content().split("\ue000")]
    return "\n".join([x for x in parts if x])

class Blocks:
    """Semantic block list parsed out of a FeedSpark deck's <body>."""
    def __init__(self):
        self.hero = None
        self.items = []          # (kind, payload) in document order
        self.close = None

def parse_deck(path, keep_checks=False):
    doc = LH.parse(path).getroot()
    body = doc.body
    # .int-note blocks are notes written to FeedSpark, not to the client -- a figure still
    # to be confirmed, where a number came from, what this deck could not reach. They are
    # dropped BEFORE anything is parsed, so there is no path by which one reaches a slide.
    # Ray asked for this after the Monsoon intro deck read as though it were addressed to
    # him rather than to the client.
    for n in body.xpath("//*[contains(concat(' ', normalize-space(@class), ' '), ' int-note ')]"):
        par = n.getparent()
        if par is not None:
            par.remove(n)
    B = Blocks()
    sc = (lambda e: txt(e, not keep_checks))

    hero = body.find(".//header[@class='hero']")
    if hero is None:
        for h in body.iter("header"):
            if has(h, "hero"): hero = h; break
    if hero is not None:
        meta = []
        for m in hero.xpath(".//div[@class='hero-meta']/div"):
            b = m.find("b")
            lab = norm(m.text or "")
            meta.append((lab, sc(b) if b is not None else ""))
        B.hero = dict(
            eyebrow=sc(hero.find(".//div[@class='eyebrow on-dark']")),
            title=sc(hero.find(".//h1")),
            lede=sc(hero.find(".//p[@class='lede']")),
            meta=meta,
        )

    close = None
    for f in body.iter("footer"):
        if has(f, "close"): close = f; break
    if close is not None:
        cts = []
        for c in close.xpath(".//div[@class='ct']"):
            b, s, a = c.find("b"), c.find("span"), c.find("a")
            cts.append((sc(b), sc(s), sc(a)))
        B.close = dict(title=sc(close.find(".//h2")),
                       eyebrow=sc(first(close, ".//div[contains(@class,'eyebrow')]")),
                       contacts=cts)

    for el in body:
        tag = el.tag
        if tag == "div" and has(el, "chapter"):
            B.items.append(("chapter", dict(
                num=sc(el.find(".//div[@class='ch-num']")),
                eyebrow=sc(first(el, ".//div[contains(@class,'eyebrow')]")),
                title=sc(el.find(".//h2")),
                sub=sc(el.find(".//p")),
            )))
        elif tag == "section":
            B.items.append(("section", parse_section(el, sc)))
    return B

def parse_section(sec, sc):
    """One <section> -> {head, blocks[]} in document order."""
    wrap = sec.find("div[@class='wrap']")
    root = wrap if wrap is not None else sec
    head = dict(eyebrow="", title="", sub="")
    blocks = []

    # NB: no id()-keyed "seen" set here. lxml builds a fresh Python proxy on each
    # element access, so a garbage-collected proxy's id() gets recycled and would
    # make unrelated elements look already-visited -- which silently dropped whole
    # sections. A plain single-visit tree walk needs no such set.
    def walk(node):
        for el in node:
            k = cls(el)
            if el.tag == "div" and "eyebrow" in k and not head["eyebrow"] and not head["title"]:
                head["eyebrow"] = sc(el); continue
            if el.tag == "h2" and "sec-title" in k and not head["title"]:
                head["title"] = sc(el); continue
            if el.tag == "p" and "sec-sub" in k and not head["sub"]:
                head["sub"] = sc(el); continue
            bs = classify(el, sc)
            if bs:
                blocks.extend(bs)
                continue
            if el.tag in ("div", "section"):
                walk(el)
    walk(root)
    return dict(head=head, blocks=blocks)

def bars_of(el):
    """Bar rows inside `el`, detected STRUCTURALLY (a .fill[data-w] under a .track).

    Class-based detection is not enough: the standard component uses .bar-row, but
    hand-built charts (e.g. the A/B test chart) use inline-styled divs with the same
    .track/.fill[data-w] skeleton. Both carry real numbers, so both must be read."""
    rows = []
    for fill in el.xpath(".//div[contains(@class,'fill')][@data-w]"):
        track = fill.getparent()
        row = track.getparent() if track is not None else None
        if row is None: continue
        spans = row.xpath(".//span")
        lab = norm(spans[0].text_content()) if spans else ""
        val = norm(spans[1].text_content()) if len(spans) > 1 else ""
        try: w = float(fill.get("data-w") or 0)
        except ValueError: w = 0
        fk = cls(fill).split()
        if lab or val:
            rows.append((lab, val, w, "green" in fk, "grey" in fk))
    return rows

def card_body(c, sc):
    paras = [sc(p) for p in c.xpath(".//p")]
    lis = [sc(li) for li in c.xpath(".//li")]
    chips = [sc(i) for i in c.xpath(".//div[contains(@style,'background:var(--paper-2)')]")]
    return "\n".join([p for p in paras if p] + [("- " + l) for l in lis] + [x for x in chips if x])

def classify(el, sc):
    """Map one element onto a LIST of renderable blocks, or None to descend into it."""
    k = cls(el).split()
    if el.tag not in ("div", "table"): return None

    if "stats" in k:
        cells = []
        for s in el.xpath("./div[contains(@class,'stat')]"):
            cells.append((sc(first(s, "div[@class='n']")), sc(first(s, "div[@class='l']"))))
        return [("stats", cells)] if cells else None

    if "agenda" in k:
        rows = []
        for r in el.xpath(".//div[@class='ag-row']"):
            rows.append((sc(first(r, "div[@class='ag-num']")),
                         sc(first(r, ".//div[@class='ag-t']")),
                         sc(first(r, ".//div[@class='ag-d']"))))
        return [("agenda", rows)] if rows else None

    if "tbl-wrap" in k or el.tag == "table":
        t = el if el.tag == "table" else first(el, ".//table")
        if t is None: return None
        heads = [sc(th) for th in t.xpath(".//thead//th")]
        rows = [[sc(td) for td in tr.xpath("./td")] for tr in t.xpath(".//tbody/tr")]
        return [("table", dict(heads=heads, rows=rows))] if rows else None

    if "sc-grid" in k:
        cells = []
        for c in el.xpath("./div[@class='sc-cell']"):
            tgt = first(c, "div[@class='sc-tgt']"); warn = first(c, "div[@class='sc-warn']")
            cells.append(dict(pct=sc(first(c, ".//div[@class='sc-pct']")),
                              label=sc(first(c, "div[@class='sc-label']")),
                              note=sc(first(c, "div[@class='sc-note']")),
                              foot=sc(tgt) if tgt is not None else sc(warn),
                              bad=warn is not None))
        return [("scorecard", cells)] if cells else None

    if "tiers" in k:
        out = []
        for t in el.xpath("./div[contains(@class,'tier')]"):
            tk = cls(t).split()
            out.append(dict(tn=sc(first(t, "div[@class='tn']")), title=sc(first(t, "h4")),
                            sub=sc(first(t, "div[@class='ts']")),
                            items=[sc(li) for li in t.xpath(".//li")],
                            here="here" in tk, done="done" in tk))
        return [("tiers", out)] if out else None

    if "road" in k:
        out = []
        for m in el.xpath("./div[contains(@class,'mo')]"):
            out.append(dict(month=sc(first(m, "div[@class='mo-m']")),
                            state=sc(first(m, ".//div[contains(@class,'mo-s')]")),
                            title=sc(first(m, "h4")),
                            items=[sc(li) for li in m.xpath(".//li")],
                            peak="peak" in cls(m).split()))
        return [("roadmap", out)] if out else None

    if "callout" in k:
        return [("callout", dict(title=sc(first(el, "h4")), body=sc(first(el, "p"))))]

    if "note" in k and "sc-note" not in k and "attr-note" not in k:
        return [("note", sc(el))]

    if "bars" in k:
        rows = bars_of(el)
        return [("bars", rows)] if rows else None

    is_grid = any(g in k for g in ("grid-2", "grid-3", "grid-4", "pipe", "flow"))
    if is_grid or "card" in k:
        cols = 2
        if "grid-3" in k or "pipe" in k: cols = 3
        if "grid-4" in k or "flow" in k: cols = 4
        kids = [c for c in el if c.tag == "div"] if is_grid else [el]
        # A card carrying a bar chart cannot render as a boxed text card without
        # losing its numbers -- decompose it into heading + full-width bars instead.
        if any(bars_of(c) for c in kids):
            out = []
            for c in kids:
                h4 = first(c, ".//h4")
                title = sc(h4) if h4 is not None else ""
                br = bars_of(c)
                if title: out.append(("subhead", title))
                if br: out.append(("bars", br))
                body = card_body(c, sc)
                if body and not br: out.append(("cards", dict(cols=1, cards=[("", body)])))
                elif body and br:
                    tail = [x for x in body.split("\n") if x.strip()]
                    if tail: out.append(("note", " ".join(tail)))
            return out or None
        cards = []
        for c in kids:
            h4 = first(c, ".//h4")
            title = sc(h4) if h4 is not None else ""
            body = card_body(c, sc)
            if not title:
                big = c.xpath("./div[contains(@style,'font-size:44px')]")
                if big:
                    title = sc(big[0])
                    lab = c.xpath("./div[contains(@style,'letter-spacing')]")
                    if lab: body = sc(lab[0])
            if title or body: cards.append((title, body))
        return [("cards", dict(cols=cols, cards=cards))] if cards else None

    # a hand-built chart panel: no component class, but real .track/.fill bars inside
    br = bars_of(el)
    if br:
        out = []
        lead = el.xpath(".//div[contains(@style,'letter-spacing')]")
        if lead:
            t = sc(lead[0])
            if t: out.append(("subhead", t))
        out.append(("bars", br))
        return out
    return None

_STATUS = [
    (r"\b(done|automated|complete|live|approved)\b", GREEN),
    (r"\b(parked|blocked|awaiting|on hold|held)\b", ORANGE_DP),
    (r"\b(in flight|in progress|semi-auto|wip|building)\b", ORANGE_DP),
    (r"^(open|scoped|to do)$", MUTED),
    (r"^high$", GREEN), (r"^medium$", ORANGE_DP), (r"^low$", MUTED),
]
def status_col(t):
    """Colour for a short status-pill cell, or None for ordinary prose."""
    v = (t or "").strip().lower()
    if not v or len(v) > 26: return None      # a sentence, not a pill
    for rx, c in _STATUS:
        if re.search(rx, v): return c
    return None
# ---------------------------------------------------------------- the template
# Every visual element -- card panels, accent bars, decorative circles, the
# gradient footer bar, the wordmark -- lives on the LAYOUTS in this template,
# not on the slides. Slides carry text only. That is what makes the output a
# real, editable PowerPoint deck rather than a picture of one: the client can
# restyle globally, re-run the theme, or drop a slide onto another layout and
# it re-flows. Never draw a rectangle here; pick a layout instead.
TPL = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates", "feedspark_deck.pptx")

A   = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
# The template's own table style, so native tables inherit its banding.
TBL_STYLE = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"

GRID  = {2: "Two-Card Grid", 3: "Three-Card Grid", 4: "Four-Card Grid",
         5: "Five-Card Grid", 6: "Six-Card Grid"}
MAXCARDS = 6

# Table body geometry, mirroring the reference deck.
TB_L, TB_W, TB_T = 0.67, 12.00, 2.02
TB_BOT   = 6.17           # must clear the Key Message strip at 6.34
HDR_H    = 0.42
ROW_MIN  = 0.42

def ph_style(layout, name):
    """(width_in, height_in, pt, bold) of a layout placeholder, read from its
    own lstStyle -- so capacity is measured against the real rendered size."""
    for p in layout.placeholders:
        if p.name != name: continue
        sz, bold = 12.0, False
        ls = p._element.find(".//" + A + "lstStyle")
        if ls is not None:
            d = ls.find(A + "lvl1pPr/" + A + "defRPr")
            if d is not None:
                if d.get("sz"): sz = int(d.get("sz")) / 100.0
                bold = d.get("b") == "1"
        return Emu(p.width).inches, Emu(p.height).inches, sz, bold
    return None

# Shrink steps tried before any text is dropped.
SCALES = (1.0, 0.94, 0.88, 0.82, 0.76, 0.70, 0.64, 0.58)
# We measure in Liberation Sans but the deck renders in Inter, so demand a
# little more room than measured rather than landing exactly on the boundary.
SAFETY = 0.96

def fit(text, geom):
    """Fit `text` to a placeholder. Returns (scale, dropped_lines, text).

    PPTX does NOT clip an over-long text frame -- it spills the text outside the
    shape, straight over the neighbouring card panel. So the failure mode to
    design against is overlap, not truncation, and the fix is to shrink rather
    than to cut: every word survives, just smaller. Only if the copy still will
    not fit at the smallest step do we drop lines, and --audit reports it."""
    if not text or geom is None: return 1.0, 0, text
    w, h, sz, bold = geom
    for s in SCALES:
        pt = sz * s
        cap = max(1, int(h / (pt * 1.28 / 72.0)))
        if len(wrap_lines(text, pt, w * SAFETY, bold)) <= cap:
            return s, 0, text
    pt = sz * SCALES[-1]
    cap = max(1, int(h / (pt * 1.28 / 72.0)))
    lines = wrap_lines(text, pt, w * SAFETY, bold)
    kept = " ".join(lines[:cap]).rstrip()
    if len(kept) > 3: kept = kept[:-1].rstrip(" ,;:.") + "…"
    return SCALES[-1], len(lines) - cap, kept

# Below this the type is too small to read on a projector; shrinking further is
# not a fit, it is a hidden failure. Content is re-laid out instead.
MIN_OK = 0.82

# Strips designed for exactly one punchy line. A paragraph poured into one of
# these shrinks to ~7pt; better to carry the lead sentence at full size.
ONE_LINERS = {"Subtitle", "Key Message", "Section Subtitle", "Attribution",
              "Date / Prepared by", "Role / Title"}

def first_sentence(t, maxlen=150):
    parts = re.split(r"(?<=[.!?])\s+", (t or "").strip())
    out = parts[0]
    for p in parts[1:]:
        if len(out) + len(p) + 1 > maxlen: break
        out += " " + p
    return out

def pick_grid(em, cards):
    """Largest card count whose panels still hold this copy at a readable size.

    Card panels are layout shapes, so a six-card layout gives six small panels
    -- pouring long copy into them is what forces 60% type. Stepping down to a
    three-card layout costs an extra slide and buys back the point size, which
    is the right trade for a client-facing deck."""
    for n in range(min(len(cards), MAXCARDS), 1, -1):
        lay = em.lay[GRID[n]]
        if all(fit(h, ph_style(lay, "Card %d Heading" % j))[0] >= MIN_OK
               and fit(b, ph_style(lay, "Card %d Body" % j))[0] >= MIN_OK
               for chunk in chunks(cards, n)
               for j, (h, b) in enumerate(chunk, start=1)):
            return n
    return 2

class Emitter:
    """Builds a deck by *populating layouts*, never by drawing shapes."""

    def __init__(self, tpl=TPL):
        if not os.path.exists(tpl):
            sys.exit("missing template: %s" % tpl)
        self.prs = Presentation(tpl)
        self.lay = {l.name: l for l in self.prs.slide_masters[0].slide_layouts}
        # Slides rename their placeholders on clone ("Text Placeholder 4"), so a
        # slide-side name lookup fails; idx survives. Map layout name -> idx here.
        self.idx = {n: {p.name: p.placeholder_format.idx for p in l.placeholders}
                    for n, l in self.lay.items()}
        self.audit = []

    def slide(self, layout):
        return self.prs.slides.add_slide(self.lay[layout])

    def put(self, slide, name, text):
        """Fill a placeholder by layout-side name. Empty -> delete the shape, so
        PowerPoint shows no 'Click to edit' prompt."""
        lname = slide.slide_layout.name
        i = self.idx[lname].get(name)
        if i is None: return
        try: shape = slide.placeholders[i]
        except KeyError: return
        lines = [l for l in str(text or "").split("\n") if l.strip()]
        if not lines:
            shape._element.getparent().remove(shape._element)
            return
        geom = ph_style(self.lay[lname], name)
        src = "\n".join(lines)
        scale, dropped, body = fit(src, geom)
        if scale < MIN_OK and name in ONE_LINERS:
            scale, dropped, body = fit(first_sentence(src), geom)
        lines = [l for l in body.split("\n") if l.strip()]
        if dropped or scale < 1.0:
            self.audit.append((len(self.prs.slides._sldIdLst), lname, name, scale, dropped))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.text = lines[0]
        for l in lines[1:]:
            tf.add_paragraph().text = l
        if scale < 1.0 and geom:
            # Set the size explicitly rather than trusting normAutofit: PowerPoint
            # recomputes fontScale on open, but LibreOffice and Google Slides do
            # not, and a deck that only fits in one renderer is not shippable.
            # Size is the sole override -- family and colour stay inherited, so
            # the layout still drives a global restyle.
            pt = Pt(round(geom[2] * scale, 1))
            for p_ in tf.paragraphs:
                for r_ in p_.runs:
                    r_.font.size = pt

    def finish(self, slide, keymsg=None):
        """Fill the Key Message strip, then drop any placeholder still empty.

        The strip's coloured panel lives on the layout, so leaving it unfilled
        renders an empty orange bar -- it must always carry text or nothing at
        all is better than a stray bar, hence the explicit delete below."""
        self.put(slide, "Key Message", keymsg or "")
        for shape in list(slide.placeholders):
            if not shape.text_frame.text.strip():
                shape._element.getparent().remove(shape._element)

    # ------------------------------------------------------------ native table
    def table(self, slide, heads, rows):
        rows = rows[:]
        ncol = max(len(heads), max((len(r) for r in rows), default=0))
        heads = (list(heads) + [""] * ncol)[:ncol]
        rows  = [(r + [""] * ncol)[:ncol] for r in rows]
        # A bar list has labels and values but no column names. Styling row 0 as
        # a header then leaves an empty coloured band that reads as a bug, so a
        # headerless table is built instead.
        hdr = any((h or "").strip() for h in heads)

        avail = TB_BOT - TB_T - (HDR_H if hdr else 0)
        body_h = min(0.75, max(ROW_MIN, avail / max(1, len(rows))))
        h = (HDR_H if hdr else 0) + body_h * len(rows)
        gf = slide.shapes.add_table(len(rows) + (1 if hdr else 0), ncol,
                                    Inches(TB_L), Inches(TB_T), Inches(TB_W), Inches(h))
        tbl = gf.table
        pr = tbl._tbl.find(A + "tblPr")
        pr.set("firstRow", "1" if hdr else "0"); pr.set("bandRow", "1")
        for e in pr.findall(A + "tableStyleId"): pr.remove(e)
        sid = pr.makeelement(A + "tableStyleId", {}); sid.text = TBL_STYLE
        pr.append(sid)

        # Columns sized by the widest cell they must carry, not evenly: a
        # "Status" column beside a prose column otherwise wraps to four lines.
        want = []
        for c in range(ncol):
            longest = max([len(heads[c])] + [len(r[c]) for r in rows] or [1])
            want.append(max(6, min(longest, 90)))
        tot = float(sum(want))
        widths = [max(0.9, TB_W * w / tot) for w in want]
        scale = TB_W / sum(widths)
        for c, w in enumerate(widths):
            tbl.columns[c].width = Emu(int(Inches(w * scale)))

        if hdr:
            tbl.rows[0].height = Inches(HDR_H)
            for i, txt_ in enumerate(heads):
                cell = tbl.cell(0, i); cell.text = txt_ or ""
                for p_ in cell.text_frame.paragraphs:
                    for r_ in p_.runs:
                        r_.font.size = Pt(13); r_.font.bold = True; r_.font.color.rgb = WHITE
        for ri, row in enumerate(rows, start=1 if hdr else 0):
            tbl.rows[ri].height = Inches(body_h)
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci); cell.text = val or ""
                col = status_col(val)
                for p_ in cell.text_frame.paragraphs:
                    for r_ in p_.runs:
                        r_.font.size = Pt(11)
                        if col is not None:
                            r_.font.bold = True; r_.font.color.rgb = col
        return gf

# ---------------------------------------------------------------- block -> slide
def chunks(seq, n):
    for i in range(0, len(seq), n):
        yield seq[i:i + n]

def balanced(seq, n):
    """Split into slides of at most n, evenly. Plain chunking turns 4 cards at
    n=3 into a 3+1 split -- a full slide followed by a near-empty one; 2+2 uses
    the same slide count and looks deliberate."""
    if len(seq) <= n: return [list(seq)]
    slides = -(-len(seq) // n)
    return list(chunks(seq, -(-len(seq) // slides)))

def grid_slide(em, title, sub, cards, key=None):
    """cards: [(heading, body)] -- laid onto the matching N-Card Grid layout.

    The card panels are layout shapes, so the layout must match the card count
    exactly or empty panels show; that is why 2..6 all exist as layouts and why
    anything longer is split rather than squeezed."""
    n = len(cards)
    if n == 1:
        s = em.slide("Title and Content")
        em.put(s, "Title", title); em.put(s, "Subtitle", sub)
        head, body = cards[0]
        em.put(s, "Content", (head + "\n" if head else "") + body)
        em.finish(s, key); return s
    s = em.slide(GRID[min(max(n, 2), MAXCARDS)])
    em.put(s, "Title", title); em.put(s, "Subtitle", sub)
    for i, (head, body) in enumerate(cards, start=1):
        em.put(s, "Card %d Heading" % i, head)
        em.put(s, "Card %d Body" % i, body)
    em.finish(s, key)
    return s

def emit_cards(em, base_title, t, sub, cards, key):
    """Emit a card list across as many slides as readability needs."""
    n = pick_grid(em, cards)
    for part in balanced(cards, n):
        grid_slide(em, t, sub, part, key)
        t, sub = (base_title + " (cont.)") if base_title else "", ""

def emit_blocks(em, head, blocks, fallback=""):
    """Flow one parsed <section> onto as many slides as its content needs."""
    # Not every <section> carries its own .sec-title -- some sit directly under a
    # chapter divider and inherit it. Without this the slide title renders as a
    # bare " (cont.)".
    title = head.get("title") or fallback or ""
    sub = head.get("sub") or ""
    first_done = False
    pending_sub = None           # a ("subhead") becomes the next slide's title
    i = 0
    while i < len(blocks):
        kind, payload = blocks[i]
        i += 1
        # A trailing note is the slide's Key Message, not a slide of its own.
        key = None
        if i < len(blocks) and blocks[i][0] == "note":
            key = blocks[i][1]; i += 1

        cont = (title + " (cont.)") if title else ""
        t = pending_sub or (title if not first_done else cont)
        s_ = sub if not first_done else ""
        pending_sub = None

        if kind == "subhead":
            pending_sub = payload
            continue

        if kind == "note":
            s = em.slide("Quote / Statement")
            em.put(s, "Quote", payload); em.put(s, "Attribution", title)
            em.finish(s); first_done = True; continue

        if kind == "callout":
            s = em.slide("Quote / Statement")
            em.put(s, "Quote", payload.get("body") or payload.get("title"))
            em.put(s, "Attribution", payload.get("title") or title)
            em.finish(s); first_done = True; continue

        if kind == "stats":
            cells = payload
            if len(cells) == 3:
                s = em.slide("Big Stats")
                em.put(s, "Title", t); em.put(s, "Subtitle", s_)
                for j, (n_, l_) in enumerate(cells, start=1):
                    em.put(s, "Stat %d Number" % j, n_)
                    em.put(s, "Stat %d Label" % j, l_)
                em.finish(s, key)
            else:
                emit_cards(em, title, t, s_, [(n_, l_) for n_, l_ in cells], key)
            first_done = True; continue

        if kind == "tiers" and len(payload) == 3:
            s = em.slide("Pricing / Tiers")
            em.put(s, "Title", t); em.put(s, "Subtitle", s_)
            for j, tr in enumerate(payload, start=1):
                em.put(s, "Tier %d Name" % j, tr.get("title") or tr.get("tn"))
                em.put(s, "Tier %d Price" % j, tr.get("tn") or "")
                em.put(s, "Tier %d Features" % j,
                       "\n".join(([tr["sub"]] if tr.get("sub") else []) + tr.get("items", [])))
            em.finish(s, key); first_done = True; continue

        if kind == "roadmap" and len(payload) == 3:
            s = em.slide("Numbered Steps")
            em.put(s, "Title", t); em.put(s, "Subtitle", s_)
            for j, m in enumerate(payload, start=1):
                em.put(s, "Step %d Number" % j, str(j))
                em.put(s, "Step %d Heading" % j, m.get("title") or m.get("month"))
                em.put(s, "Step %d Detail" % j,
                       "\n".join(([m["month"]] if m.get("month") else []) + m.get("items", [])))
            em.finish(s, key); first_done = True; continue

        if kind in ("tiers", "roadmap"):
            cards = []
            for x in payload:
                head_ = x.get("title") or x.get("month") or x.get("tn") or ""
                body_ = "\n".join(([x.get("sub") or x.get("month") or ""] if (x.get("sub") or x.get("month")) else [])
                                  + x.get("items", []))
                cards.append((head_, body_))
            emit_cards(em, title, t, s_, cards, key)
            first_done = True; continue

        if kind == "scorecard":
            cards = [((c["pct"] + "  " + c["label"]).strip(),
                      "\n".join([x for x in (c.get("note"), c.get("foot")) if x]))
                     for c in payload]
            emit_cards(em, title, t, s_, cards, key)
            first_done = True; continue

        if kind == "cards":
            cards = [(h or "", b or "") for h, b in payload["cards"]]
            emit_cards(em, title, t, s_, cards, key)
            first_done = True; continue

        if kind == "agenda":
            rows = ["%s  %s — %s" % (n_, ti, de) for n_, ti, de in payload]
            half = (len(rows) + 1) // 2
            s = em.slide("Two Content")
            em.put(s, "Title", t); em.put(s, "Subtitle", s_)
            em.put(s, "Left Content", "\n".join(rows[:half]))
            em.put(s, "Right Content", "\n".join(rows[half:]))
            em.finish(s, key); first_done = True; continue

        if kind == "bars":
            heads = ["", ""]
            rws = [[lab, val] for lab, val, _w, _g, _y in payload]
            for part in chunks(rws, 8):
                s = em.slide("Table")
                em.put(s, "Title", t); em.put(s, "Subtitle", s_)
                em.table(s, heads, part); em.finish(s, key)
                t, s_ = cont, ""
            first_done = True; continue

        if kind == "table":
            heads, rws = payload["heads"], payload["rows"]
            per = max(1, int((TB_BOT - TB_T - HDR_H) / ROW_MIN))
            for part in chunks(rws, per):
                s = em.slide("Table")
                em.put(s, "Title", t); em.put(s, "Subtitle", s_)
                em.table(s, heads, part); em.finish(s, key)
                t, s_ = cont, ""
            first_done = True; continue

def build(src, out, tpl=TPL, keep_checks=False):
    B = parse_deck(src, keep_checks)
    em = Emitter(tpl)

    if B.hero:
        s = em.slide("Title Slide")
        em.put(s, "Title", B.hero["title"])
        lede = B.hero.get("lede") or ""
        pull = len(lede) > 90
        # A short lede is a subtitle. A long one is a pull-quote, and squeezing it
        # into a one-line strip shrinks it to ~7pt -- it earns its own slide.
        em.put(s, "Subtitle", B.hero["eyebrow"] if pull else (lede or B.hero["eyebrow"]))
        em.put(s, "Date / Prepared by",
               "  ·  ".join("%s %s" % (l, v) for l, v in B.hero["meta"] if v))
        em.finish(s)
        if pull:
            q = em.slide("Quote / Statement")
            em.put(q, "Quote", lede)
            em.finish(q)          # no attribution invented for an unsourced lede

    last_ch = ""
    for kind, payload in B.items:
        if kind == "chapter":
            last_ch = payload["title"]
            s = em.slide("Section Marker")
            em.put(s, "Section Label", payload.get("eyebrow") or payload.get("num"))
            em.put(s, "Section Title", payload["title"])
            em.put(s, "Section Subtitle", payload.get("sub"))
            em.finish(s)
        else:
            emit_blocks(em, payload["head"], payload["blocks"], last_ch)

    if B.close:
        s = em.slide("Closing")
        em.put(s, "Title", B.close.get("title") or "Questions?")
        cts = B.close.get("contacts") or []
        if cts:
            n_, r_, a_ = cts[0]
            em.put(s, "Name", n_); em.put(s, "Role / Title", r_); em.put(s, "Contact details", a_)
        if len(cts) > 1:
            em.put(s, "Second-Column Heading", "Also on the account")
            em.put(s, "Second-Column Content",
                   "\n".join("%s — %s" % (n_, a_ or r_) for n_, r_, a_ in cts[1:]))
        em.finish(s)

    em.prs.save(out)
    return em

def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("src"); ap.add_argument("out")
    ap.add_argument("--template", default=TPL)
    ap.add_argument("--keep-checks", action="store_true",
                    help="keep the `?` data-check badges (dropped by default: the web "
                         "deck has a toggle to hide them, a .pptx does not)")
    ap.add_argument("--audit", action="store_true",
                    help="report layout usage and any text trimmed to fit")
    a = ap.parse_args()
    em = build(a.src, a.out, a.template, a.keep_checks)
    n = len(em.prs.slides._sldIdLst)
    print("wrote %s -- %d slides" % (a.out, n))
    if a.audit:
        used = {}
        for s in em.prs.slides:
            used[s.slide_layout.name] = used.get(s.slide_layout.name, 0) + 1
        print("\nlayouts used:")
        for k, v in sorted(used.items(), key=lambda x: -x[1]):
            print("  %-26s %d" % (k, v))
        shrunk = [a for a in em.audit if a[4] == 0]
        cut    = [a for a in em.audit if a[4] > 0]
        print("\nshrunk to fit: %d   |   still over capacity: %d" % (len(shrunk), len(cut)))
        for sn, lay, name, scale, dropped in cut:
            print("  DROPPED slide %-3d %-20s %-22s -%d line(s)" % (sn, lay, name, dropped))
        tight = sorted(shrunk, key=lambda a: a[3])[:8]
        if tight:
            print("\n  tightest fits:")
            for sn, lay, name, scale, _d in tight:
                print("    slide %-3d %-20s %-22s %d%%" % (sn, lay, name, round(scale * 100)))

if __name__ == "__main__":
    main()
