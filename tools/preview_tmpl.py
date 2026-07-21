#!/usr/bin/env python3
"""Template-aware pptx -> PNG previewer: composites layout decorative shapes
(cards, accent bars, gradient bars) UNDER each slide's own placeholder content,
since real PowerPoint templates keep repeating chrome on the layout, not the slide.
Falls back to the layout placeholder's formatting when a slide placeholder run
has no explicit color (common for templated decks relying on inherited style)."""
import re, sys
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

SC = 120
EMU = 914400
LIB = "/usr/share/fonts/truetype/liberation/"
DEJ = "/usr/share/fonts/truetype/dejavu/"

def px(emu): return int(round(emu / EMU * SC))
def rgb(c):
    if c is None: return None
    h = str(c)
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

_fc = {}
def font(name, size_px, bold):
    size_px = max(size_px, 6)
    key=(name,size_px,bold)
    if key in _fc: return _fc[key]
    n=(name or "").lower()
    if "consol" in n or "mono" in n:
        p = DEJ+("DejaVuSansMono-Bold.ttf" if bold else "DejaVuSansMono.ttf")
        f=ImageFont.truetype(p,size_px)
    else:
        p = LIB+("LiberationSans-Bold.ttf" if bold else "LiberationSans-Regular.ttf")
        f=ImageFont.truetype(p,size_px)
    _fc[key]=f; return f

def solid_fill(obj):
    try:
        f=obj.fill
        if f.type==1: return rgb(f.fore_color.rgb)
        if f.type==3:  # gradient: approximate with first stop
            try:
                stops = f.gradient_stops
                return rgb(stops[0].color.rgb)
            except Exception:
                return (200,200,200)
    except Exception: pass
    return None

def gradient_fill(obj):
    try:
        f=obj.fill
        if f.type==3:
            stops = list(f.gradient_stops)
            cols = [rgb(s.color.rgb) for s in stops]
            return cols
    except Exception: pass
    return None

def line_spec(shape):
    try:
        ln=shape.line
        if ln.fill.type==1:
            w = ln.width.pt if ln.width else 1.0
            return rgb(ln.color.rgb), w
    except Exception: pass
    return None,0

def slide_bg(slide, layout):
    for src in (slide, layout):
        try:
            if src.background.fill.type==1:
                return rgb(src.background.fill.fore_color.rgb)
        except Exception: pass
    return (255,255,255)

def defrpr_default(text_frame, level=0):
    """Paragraph-level default run formatting (a:lstStyle/a:lvlNpPr/a:defRPr) —
    the style path templated decks use for placeholders with no literal runs."""
    try:
        body = text_frame._txBody
        lst = body.find(qn('a:lstStyle'))
        if lst is None: return (None, None, None, None)
        lvl_tag = f'a:lvl{level+1}pPr'
        lvlPr = lst.find(qn(lvl_tag))
        if lvlPr is None: lvlPr = lst.find(qn('a:lvl1pPr'))
        if lvlPr is None: return (None, None, None, None)
        defRPr = lvlPr.find(qn('a:defRPr'))
        if defRPr is None: return (None, None, None, None)
        sz = defRPr.get('sz')
        sz = int(sz)/100 if sz else None
        b = defRPr.get('b')
        b = (b == '1') if b is not None else None
        fill = defRPr.find(qn('a:solidFill'))
        col = None
        if fill is not None:
            clr = fill.find(qn('a:srgbClr'))
            if clr is not None: col = rgb(clr.get('val'))
        latin = defRPr.find(qn('a:latin'))
        name = latin.get('typeface') if latin is not None else None
        return (sz, b, col, name)
    except Exception:
        return (None, None, None, None)

def para_runs(p, fallback_runs=None, defrpr=None):
    out=[]
    fb = fallback_runs or []
    dsz, db, dcol, dname = defrpr or (None, None, None, None)
    for i, r in enumerate(p.runs):
        sz = r.font.size.pt if r.font.size else None
        b  = r.font.bold
        col= rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else None
        if i < len(fb):
            fsz, fb_, fcol, fname = fb[i]
            if sz is None: sz = fsz
            if b is None: b = fb_
            if col is None: col = fcol
            name = r.font.name or fname
        else:
            name = r.font.name
        # deepest fallback: paragraph-level lstStyle defRPr (title/subtitle placeholders
        # in templated decks usually carry NO literal runs at all on the layout)
        if sz is None: sz = dsz
        if b is None: b = db
        if col is None: col = dcol
        if name is None: name = dname
        out.append((r.text, name, sz or 12, bool(b), col or (20,29,43)))
    return out

def collect_fallback(p):
    out = []
    for r in p.runs:
        sz = r.font.size.pt if r.font.size else 12
        b = bool(r.font.bold)
        col = rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else None
        out.append((sz, b, col, r.font.name))
    return out

def para_indent(p):
    pPr = p._p.find(qn('a:pPr'))
    if pPr is None: return 0.0, 0.0
    marL = int(pPr.get('marL', 0)) / EMU * SC
    indent = int(pPr.get('indent', 0)) / EMU * SC
    return marL, indent

def para_align(p, layout_p):
    if p.alignment is not None: return p.alignment
    if layout_p is not None and layout_p.alignment is not None: return layout_p.alignment
    return PP_ALIGN.LEFT

def wrap(draw, runs, box_w, first_extra=0.0):
    tokens=[]
    for text,name,sz,b,col in runs:
        for part in re.split(r'(\s+)', text):
            if part=='' : continue
            tokens.append((part,name,sz,b,col))
    lines=[[]]; cur=0.0; avail=box_w+first_extra
    for (wd,name,sz,b,col) in tokens:
        f=font(name,int(sz*SC/72),b)
        wdt=draw.textlength(wd,font=f)
        isspace = wd.strip()==''
        if not isspace and cur+wdt>avail and lines[-1]:
            lines.append([]); cur=0.0; avail=box_w
        if isspace and not lines[-1]:
            continue
        lines[-1].append((wd,f,col,wdt,sz)); cur+=wdt
    return [ln for ln in lines if ln] or [[]]

def draw_textframe(draw, shape, warn, layout_shape=None):
    tf=shape.text_frame
    x=px(shape.left); y=px(shape.top); w=px(shape.width); h=px(shape.height)
    box_w=max(4,w-4)
    layout_paras = list(layout_shape.text_frame.paragraphs) if (layout_shape is not None and layout_shape.has_text_frame) else []
    defrpr = defrpr_default(layout_shape.text_frame) if layout_shape is not None and layout_shape.has_text_frame else None
    if defrpr is None or all(v is None for v in defrpr):
        defrpr = defrpr_default(tf)
    paras=[]
    total=0.0
    for pi, p in enumerate(tf.paragraphs):
        lp = layout_paras[pi] if pi < len(layout_paras) else (layout_paras[0] if layout_paras else None)
        fb = collect_fallback(lp) if lp is not None else None
        runs=para_runs(p, fb, defrpr)
        marL, indent = para_indent(p)
        if not any(t.strip() for t,_,_,_,_ in runs):
            sz = runs[0][2] if runs else 12
            paras.append((None,p,sz,0,0,None)); total+= sz*SC/72*1.1
            continue
        lines=wrap(draw,runs,box_w-marL,first_extra=-indent)
        ls = p.line_spacing if p.line_spacing else 1.15
        sa = p.space_after.pt if p.space_after else 0
        pheight=0.0
        for ln in lines:
            lh=max((sz for *_,sz in ln), default=12)*SC/72*ls
            pheight+=lh
        pheight+=sa*SC/72
        align = para_align(p, lp)
        paras.append((lines,p,None,marL,indent,align)); total+=pheight
    anchor=tf.vertical_anchor or (layout_shape.text_frame.vertical_anchor if layout_shape is not None and layout_shape.has_text_frame else None)
    if anchor==MSO_ANCHOR.MIDDLE: cy=y+max(0,(h-total)/2)
    elif anchor==MSO_ANCHOR.BOTTOM: cy=y+max(0,h-total)
    else: cy=y
    if total>h+4: warn.append((shape,total,h))
    for item in paras:
        lines,p,blanksz,marL,indent,align=item
        if blanksz is not None:
            cy+=blanksz*SC/72*1.1; continue
        ls=p.line_spacing if p.line_spacing else 1.15
        sa=p.space_after.pt if p.space_after else 0
        for li,ln in enumerate(lines):
            lh=max((sz for *_,sz in ln),default=12)*SC/72*ls
            lw=sum(wdt for _,_,_,wdt,_ in ln)
            off = (marL+indent) if li==0 else marL
            if align==PP_ALIGN.CENTER: sx=x+(w-lw)/2
            elif align==PP_ALIGN.RIGHT: sx=x+w-lw
            else: sx=x+off
            asc = max((sz for *_,sz in ln),default=12)*SC/72
            ty=cy+(lh-asc)/2
            for (wd,f,col,wdt,sz) in ln:
                draw.text((sx,ty),wd,font=f,fill=col or (20,29,43))
                sx+=wdt
            cy+=lh
        cy+=sa*SC/72

def draw_table(d, shape):
    tbl = shape.table
    x0 = px(shape.left); y0 = px(shape.top)
    col_w = [px(c.width) for c in tbl.columns]
    row_h = [px(r.height) for r in tbl.rows]
    from pptx.enum.text import PP_ALIGN as _A
    y = y0
    for ri, row in enumerate(tbl.rows):
        x = x0
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci]; rh = row_h[ri]
            fillc = None
            try:
                if cell.fill.type == 1: fillc = rgb(cell.fill.fore_color.rgb)
            except Exception: pass
            if fillc: d.rectangle([x, y, x+cw, y+rh], fill=fillc)
            d.rectangle([x, y, x+cw, y+rh], outline=(226,232,240), width=1)
            para = cell.text_frame.paragraphs[0]
            runs = para_runs(para)
            if runs and any(t.strip() for t,_,_,_,_ in runs):
                txt = "".join(t for t,_,_,_,_ in runs)
                _, name, sz, b, col = runs[0]
                f = font(name, int(sz*SC/72), b)
                pad = int(0.06*SC)
                tw = d.textlength(txt, font=f)
                if para.alignment == _A.CENTER: tx = x + (cw-tw)/2
                elif para.alignment == _A.RIGHT: tx = x + cw - tw - pad
                else: tx = x + pad
                ty = y + (rh - sz*SC/72)/2
                d.text((tx, ty), txt, font=f, fill=col or (20,29,43))
            x += cw
        y += row_h[ri]

def draw_shape(d, shape, warn, layout_shape=None):
    try:
        x=px(shape.left); y=px(shape.top); w=px(shape.width); h=px(shape.height)
    except Exception:
        return
    if w<=0 or h<=0: return
    if getattr(shape, "has_table", False) and shape.has_table:
        draw_table(d, shape); return
    cls=shape.__class__.__name__
    if cls=="Connector":
        col,lw=line_spec(shape)
        if col: d.line([x,y,x+w,y+h],fill=col,width=max(1,int(lw*SC/72)))
        return
    grad = gradient_fill(shape)
    fill=solid_fill(shape); lcol,lw=line_spec(shape)
    auto=None
    try: auto=shape.auto_shape_type
    except Exception: auto=None
    an=str(auto)
    lwpx=max(1,int(lw*SC/72)) if lcol else 0
    # only ovals need special handling: a naive rectangular gradient paste on a round
    # "glow" shape looks like a solid block, whereas thin gradient bars/lines render fine
    is_decorative_glow = grad and "OVAL" in an
    if grad and len(grad) >= 2 and not is_decorative_glow:
        c1, c2 = grad[0], grad[-1]
        band = Image.new("RGB", (w, h))
        bd = ImageDraw.Draw(band)
        for i in range(w):
            t = i / max(1, w-1)
            col = tuple(int(c1[k]+(c2[k]-c1[k])*t) for k in range(3))
            bd.line([(i,0),(i,h)], fill=col)
        d._image.paste(band, (x,y))
    elif grad and is_decorative_glow:
        pass  # skip large/oval gradient "glow" flourishes — can't render alpha-fade faithfully, so don't misrepresent as a solid block
    elif fill or lcol:
        if "OVAL" in an:
            d.ellipse([x,y,x+w,y+h],fill=fill,outline=lcol,width=lwpx or (1 if lcol else 0))
        elif "ROUNDED" in an:
            try: adj=shape.adjustments[0]
            except Exception: adj=0.1
            rad=int(min(w,h)*min(max(adj,0),0.5))
            rad=max(0,min(rad,min(w,h)//2-1))
            if rad>0:
                d.rounded_rectangle([x,y,x+w-1,y+h-1],radius=rad,fill=fill,outline=lcol,width=lwpx or (1 if lcol else 0))
            else:
                d.rectangle([x,y,x+w-1,y+h-1],fill=fill,outline=lcol,width=lwpx or (1 if lcol else 0))
        elif shape.shape_type != 13:  # not PICTURE
            d.rectangle([x,y,x+w-1,y+h-1],fill=fill,outline=lcol,width=lwpx or (1 if lcol else 0))
    if shape.has_text_frame and shape.text_frame.text.strip():
        draw_textframe(d,shape,warn, layout_shape=layout_shape)

def render(path, outbase):
    prs=Presentation(path)
    W=px(prs.slide_width); Hh=px(prs.slide_height)
    warns_all=[]
    n = len(prs.slides._sldIdLst)
    for i,slide in enumerate(prs.slides):
        layout = slide.slide_layout
        img=Image.new("RGB",(W,Hh),slide_bg(slide, layout))
        d=ImageDraw.Draw(img)
        d._image = img
        warn=[]

        slide_ph_idx = set()
        for sh in slide.placeholders:
            slide_ph_idx.add(sh.placeholder_format.idx)

        # 1) layout decorative + unoverridden placeholder shapes (background chrome)
        for lsh in layout.shapes:
            if lsh.is_placeholder and lsh.placeholder_format.idx in slide_ph_idx:
                continue  # slide provides its own version, skip layout ghost text
            if lsh.is_placeholder and lsh.placeholder_format.type == 13:  # SLIDE_NUMBER
                continue
            draw_shape(d, lsh, warn)

        # 2) slide's own shapes on top, matched to layout counterpart for style fallback
        layout_by_idx = {sh.placeholder_format.idx: sh for sh in layout.placeholders} if hasattr(layout, 'placeholders') else {}
        for sh in slide.shapes:
            lsh = None
            if sh.is_placeholder:
                lsh = layout_by_idx.get(sh.placeholder_format.idx)
            draw_shape(d, sh, warn, layout_shape=lsh)

        img.save(f"{outbase}_{i+1}.png")
        for (sh,t,hh) in warn:
            label = sh.text_frame.text[:48] if sh.has_text_frame else sh.name
            warns_all.append((i+1, label, round(t), hh))
    print("rendered", n, "slides")
    if warns_all:
        print("OVERFLOW WARNINGS (slide, text, text_h_px, box_h_px):")
        for wv in warns_all: print("  ", wv)
    else:
        print("no overflow warnings")

if __name__=="__main__":
    render(sys.argv[1], sys.argv[2])
