#!/usr/bin/env python3
"""Faithful-ish pptx -> PNG previewer for QA: reads real shape geometry/text
from the .pptx and draws with Pillow using Liberation Sans (Arial-metric-compatible),
so word-wrap / overflow closely matches PowerPoint."""
import re, sys
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

SC = 120  # px per inch
EMU = 914400
LIB = "/usr/share/fonts/truetype/liberation/"
DEJ = "/usr/share/fonts/truetype/dejavu/"

def px(emu): return int(round(emu / EMU * SC))
def rgb(c):
    if c is None: return None
    h = str(c)
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

_fc = {}
def font(name, size_px, bold):
    key=(name,size_px,bold)
    if key in _fc: return _fc[key]
    n=(name or "").lower()
    if "consol" in n or "mono" in n:
        p = DEJ+("DejaVuSansMono-Bold.ttf" if bold else "DejaVuSansMono.ttf")
        try: f=ImageFont.truetype(p,size_px)
        except Exception: f=ImageFont.truetype(DEJ+"DejaVuSansMono.ttf",size_px)
    else:
        p = LIB+("LiberationSans-Bold.ttf" if bold else "LiberationSans-Regular.ttf")
        f=ImageFont.truetype(p,size_px)
    _fc[key]=f; return f

def solid_fill(obj):
    try:
        f=obj.fill
        if f.type==1: return rgb(f.fore_color.rgb)
    except Exception: pass
    return None
def line_spec(shape):
    try:
        ln=shape.line
        if ln.fill.type==1:
            w = ln.width.pt if ln.width else 1.0
            return rgb(ln.color.rgb), w
    except Exception: pass
    return None,0

def slide_bg(slide):
    try:
        if slide.background.fill.type==1:
            return rgb(slide.background.fill.fore_color.rgb)
    except Exception: pass
    return (255,255,255)

def para_runs(p):
    out=[]
    for r in p.runs:
        sz = r.font.size.pt if r.font.size else 12
        b  = bool(r.font.bold)
        col= rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else (20,29,43)
        out.append((r.text, r.font.name, sz, b, col))
    return out

def para_indent(p):
    pPr = p._p.find(qn('a:pPr'))
    if pPr is None: return 0.0, 0.0
    marL = int(pPr.get('marL', 0)) / EMU * SC
    indent = int(pPr.get('indent', 0)) / EMU * SC
    return marL, indent

def wrap(draw, runs, box_w, first_extra=0.0):
    tokens=[]
    for text,name,sz,b,col in runs:
        for part in re.split(r'(\s+)', text):
            if part=='' : continue
            tokens.append((part,name,sz,b,col))
    lines=[[]]; cur=0.0; avail=box_w+first_extra
    for (wd,name,sz,b,col) in tokens:
        f=font(name,int(sz*SC/72),b)
        wdt=draw.textlength(wd,font=f)
        isspace = wd.strip()==''
        if not isspace and cur+wdt>avail and lines[-1]:
            lines.append([]); cur=0.0; avail=box_w
        if isspace and not lines[-1]:
            continue
        lines[-1].append((wd,f,col,wdt,sz)); cur+=wdt
    return [ln for ln in lines if ln] or [[]]

def draw_textframe(draw, shape, warn):
    tf=shape.text_frame
    x=px(shape.left); y=px(shape.top); w=px(shape.width); h=px(shape.height)
    box_w=max(4,w-2)
    paras=[]
    total=0.0
    for p in tf.paragraphs:
        runs=para_runs(p)
        marL, indent = para_indent(p)
        if not any(t.strip() for t,_,_,_,_ in runs):
            # blank paragraph -> still take vertical space of its font
            sz = runs[0][2] if runs else 12
            paras.append((None,p,sz,0,0)); total+= sz*SC/72*1.1
            continue
        lines=wrap(draw,runs,box_w-marL,first_extra=-indent)
        ls = p.line_spacing if p.line_spacing else 1.12
        sa = p.space_after.pt if p.space_after else 0
        pheight=0.0
        for ln in lines:
            lh=max((sz for *_,sz in ln), default=12)*SC/72*ls
            pheight+=lh
        pheight+=sa*SC/72
        paras.append((lines,p,None,marL,indent)); total+=pheight
    anchor=tf.vertical_anchor
    if anchor==MSO_ANCHOR.MIDDLE: cy=y+(h-total)/2
    elif anchor==MSO_ANCHOR.BOTTOM: cy=y+h-total
    else: cy=y
    if total>h+2: warn.append((shape,total,h))
    for item in paras:
        lines,p,blanksz,marL,indent=item
        if blanksz is not None:
            cy+=blanksz*SC/72*1.1; continue
        ls=p.line_spacing if p.line_spacing else 1.12
        sa=p.space_after.pt if p.space_after else 0
        align=p.alignment
        for li,ln in enumerate(lines):
            lh=max((sz for *_,sz in ln),default=12)*SC/72*ls
            lw=sum(wdt for _,_,_,wdt,_ in ln)
            off = (marL+indent) if li==0 else marL
            if align==PP_ALIGN.CENTER: sx=x+(w-lw)/2
            elif align==PP_ALIGN.RIGHT: sx=x+w-lw
            else: sx=x+off
            # baseline: place text top at cy + (lh - glyph)/2 approx
            asc = max((sz for *_,sz in ln),default=12)*SC/72
            ty=cy+(lh-asc)/2
            for (wd,f,col,wdt,sz) in ln:
                draw.text((sx,ty),wd,font=f,fill=col or (20,29,43))
                sx+=wdt
            cy+=lh
        cy+=sa*SC/72

def render(path, outbase):
    prs=Presentation(path)
    W=px(prs.slide_width); Hh=px(prs.slide_height)
    warns_all=[]
    for i,slide in enumerate(prs.slides):
        img=Image.new("RGB",(W,Hh),slide_bg(slide))
        d=ImageDraw.Draw(img)
        warn=[]
        for shape in slide.shapes:
            cls=shape.__class__.__name__
            x=px(shape.left); y=px(shape.top); w=px(shape.width); h=px(shape.height)
            if cls=="Connector":
                col,lw=line_spec(shape)
                if col: d.line([x,y,x+w,y+h],fill=col,width=max(1,int(lw*SC/72)))
                continue
            fill=solid_fill(shape); lcol,lw=line_spec(shape)
            auto=None
            try: auto=shape.auto_shape_type
            except Exception: auto=None
            an=str(auto)
            lwpx=max(1,int(lw*SC/72)) if lcol else 0
            if fill or lcol:
                if "OVAL" in an:
                    d.ellipse([x,y,x+w,y+h],fill=fill,outline=lcol,width=lwpx or 1 if lcol else 0)
                elif "ROUNDED" in an:
                    try: adj=shape.adjustments[0]
                    except Exception: adj=0.1
                    rad=int(min(w,h)*min(max(adj,0),0.5))
                    rad=max(0,min(rad,min(w,h)//2-1))
                    if rad>0:
                        d.rounded_rectangle([x,y,x+w-1,y+h-1],radius=rad,fill=fill,outline=lcol,width=lwpx or 1 if lcol else 0)
                    else:
                        d.rectangle([x,y,x+w-1,y+h-1],fill=fill,outline=lcol,width=lwpx or 1 if lcol else 0)
                else:
                    d.rectangle([x,y,x+w-1,y+h-1],fill=fill,outline=lcol,width=lwpx or 1 if lcol else 0)
            if shape.has_text_frame and shape.text_frame.text.strip():
                draw_textframe(d,shape,warn)
        img.save(f"{outbase}_{i+1}.png")
        for (sh,t,hh) in warn:
            warns_all.append((i+1, sh.text_frame.text[:48], round(t), hh))
    print("rendered", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
    if warns_all:
        print("OVERFLOW WARNINGS (slide, text, text_h_px, box_h_px):")
        for wv in warns_all: print("  ", wv)
    else:
        print("no overflow warnings")

if __name__=="__main__":
    render(sys.argv[1], sys.argv[2])
